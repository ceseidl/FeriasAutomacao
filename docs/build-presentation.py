"""
Gera docs/Apresentacao-FeriasAutomacao.pptx -- apresentacao executiva
do projeto Ferias Automacao para chefia.

Uso:
    python docs/build-presentation.py
"""
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt

# ---------------------------------------------------------------------------
# Identidade visual
# ---------------------------------------------------------------------------
PRIMARY   = RGBColor(0x1F, 0x49, 0x7D)   # azul corporativo
ACCENT    = RGBColor(0x2E, 0x75, 0xB6)   # azul claro
INK       = RGBColor(0x26, 0x26, 0x26)   # texto principal
MUTED     = RGBColor(0x59, 0x59, 0x59)   # texto secundario
DIVIDER   = RGBColor(0xBF, 0xBF, 0xBF)   # cinza linha
SUCCESS   = RGBColor(0x1E, 0x7E, 0x34)   # verde status OK
WARN      = RGBColor(0xC9, 0x73, 0x0F)   # laranja status atencao
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_BG  = RGBColor(0xF2, 0xF2, 0xF2)

FONT = "Arial"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(0.95)


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def add_textbox(slide, left, top, width, height, text, *,
                font_size=18, bold=False, color=INK, align=PP_ALIGN.LEFT,
                anchor=MSO_ANCHOR.TOP, font=FONT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.05)
    tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.02)
    tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return box


def add_filled_rect(slide, left, top, width, height, fill_color, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = DIVIDER
        shape.line.width = Pt(0.5)
    shape.shadow.inherit = False
    return shape


def add_header(slide, title, subtitle=None):
    add_filled_rect(slide, Inches(0), Inches(0), SLIDE_W, HEADER_H, PRIMARY)
    add_textbox(slide, Inches(0.5), Inches(0.12), Inches(12), Inches(0.55),
                title, font_size=26, bold=True, color=WHITE)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(0.6), Inches(12), Inches(0.35),
                    subtitle, font_size=13, color=WHITE)


def add_footer(slide, slide_no, total):
    add_filled_rect(slide, Inches(0), Inches(7.2), SLIDE_W, Inches(0.04), ACCENT)
    add_textbox(slide, Inches(0.5), Inches(7.25), Inches(8), Inches(0.25),
                "Ferias Automacao  |  Apresentacao executiva",
                font_size=9, color=MUTED)
    add_textbox(slide, Inches(11.7), Inches(7.25), Inches(1.2), Inches(0.25),
                f"{slide_no} / {total}",
                font_size=9, color=MUTED, align=PP_ALIGN.RIGHT)


def add_bullets(slide, left, top, width, height, bullets, *,
                font_size=16, color=INK, bullet_color=ACCENT, line_spacing=1.25):
    """bullets: list of (text, bold) tuples or just strings."""
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.05)

    for i, item in enumerate(bullets):
        if isinstance(item, tuple):
            text, bold = item
        else:
            text, bold = item, False
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        # bullet marker
        bullet = p.add_run()
        bullet.text = "\u25A0  "
        bullet.font.name = FONT
        bullet.font.size = Pt(font_size)
        bullet.font.color.rgb = bullet_color
        bullet.font.bold = True
        # actual text
        run = p.add_run()
        run.text = text
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
    return box


def add_kv_block(slide, left, top, width, height, rows, *,
                 key_width_in=2.6, font_size=14):
    """Two-column 'label / value' table-like block, rendered as boxes."""
    rect = add_filled_rect(slide, left, top, width, height, LIGHT_BG, line=True)
    row_h = (height - Inches(0.2)) / max(len(rows), 1)
    cur_top = top + Inches(0.1)
    for label, value in rows:
        add_textbox(slide, left + Inches(0.2), cur_top,
                    Inches(key_width_in), row_h,
                    label, font_size=font_size, bold=True, color=PRIMARY,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, left + Inches(0.2 + key_width_in), cur_top,
                    width - Inches(0.4 + key_width_in), row_h,
                    value, font_size=font_size, color=INK,
                    anchor=MSO_ANCHOR.MIDDLE)
        cur_top += row_h


def add_pipeline_step(slide, left, top, width, height, title, body, color=ACCENT):
    add_filled_rect(slide, left, top, width, height, color)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.08),
                width - Inches(0.2), Inches(0.45),
                title, font_size=14, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, left + Inches(0.1), top + Inches(0.55),
                width - Inches(0.2), height - Inches(0.6),
                body, font_size=11, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)


def add_arrow(slide, left, top, width, height):
    arrow = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, left, top, width, height)
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = MUTED
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    return arrow


# ---------------------------------------------------------------------------
# Content
# ---------------------------------------------------------------------------
TODAY = date.today().strftime("%d/%m/%Y")
AUTHOR = "Carlos Seidl"
VERSION = "1.0.0"

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank = prs.slide_layouts[6]
TOTAL_SLIDES = 16


# ---------------------------------------------------------------------------
# Slide 1 -- Capa
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_filled_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY)
add_filled_rect(s, Inches(0), Inches(5.4), SLIDE_W, Inches(0.06), ACCENT)
add_textbox(s, Inches(0.8), Inches(2.3), Inches(11.5), Inches(1.0),
            "Ferias Automacao", font_size=54, bold=True, color=WHITE)
add_textbox(s, Inches(0.8), Inches(3.4), Inches(11.5), Inches(0.6),
            "Geracao automatica do Planejamento de Ferias",
            font_size=22, color=WHITE)
add_textbox(s, Inches(0.8), Inches(4.05), Inches(11.5), Inches(0.5),
            "Apresentacao executiva  |  Status do projeto e proximos passos",
            font_size=14, color=WHITE)
add_textbox(s, Inches(0.8), Inches(5.7), Inches(11.5), Inches(0.4),
            f"Autor: {AUTHOR}", font_size=14, color=WHITE)
add_textbox(s, Inches(0.8), Inches(6.1), Inches(11.5), Inches(0.4),
            f"Data: {TODAY}", font_size=14, color=WHITE)
add_textbox(s, Inches(0.8), Inches(6.5), Inches(11.5), Inches(0.4),
            f"Versao do app: {VERSION}", font_size=14, color=WHITE)


# ---------------------------------------------------------------------------
# Slide 2 -- Contexto / problema
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_header(s, "Contexto e problema",
           "O que existia antes do app")
add_bullets(s, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), [
    ("Planejamento de ferias era mantido em planilha solta, sem padronizacao visual.", True),
    "Lideres precisavam abrir o Excel e interpretar tabelas para entender o panorama do mes.",
    "Compartilhamento exigia print, copia manual ou envio do .xlsx -- sem visual executivo.",
    "Inexistia um cronograma visual (Gantt) consolidado por squad e por mes.",
    "Cada pessoa formatava o relatorio do seu jeito -- inconsistencia entre squads.",
    ("Resultado: tempo perdido formatando, retrabalho e baixa visibilidade para a chefia.", True),
])
add_footer(s, 2, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Slide 3 -- Solucao em alto nivel
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_header(s, "Solucao",
           "Um app desktop que transforma a planilha em relatorio executivo em segundos")
add_bullets(s, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), [
    ("Entrada unica: planilha-template oficial Ferias-template.xlsx", True),
    "Saida automatica: HTML estilizado, Markdown, copia da planilha e PDF (opcional).",
    "Dashboard mensal com qtd. de pessoas, squads envolvidos e status geral.",
    "Cronograma detalhado linha-a-linha por colaborador.",
    "Gantt dinamico (Mermaid) por mes -- visualizacao executiva imediata.",
    "Rodape de autoria com nome do gerador e data/hora -- rastreabilidade.",
    ("Foco: zero configuracao no PC do usuario -- duplo-clique e gerar.", True),
])
add_footer(s, 3, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Slide 4 -- Stack tecnica
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_header(s, "Stack tecnica",
           "Componentes nativos do Windows + ferramentas open-source consolidadas")
add_kv_block(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(5.4), [
    ("Sistema",        "Windows 10 / 11"),
    ("Runtime",        "PowerShell 5.1+ (ja vem instalado no Windows)"),
    ("Conversor",      "Pandoc -- gera HTML standalone com CSS embutido"),
    ("Visual Gantt",   "Mermaid (renderizado no navegador via CDN)"),
    ("Leitura xlsx",   "Modulo PowerShell ImportExcel (auto-instalado)"),
    ("Interface",      "Windows Forms (WinForms) -- nativo do .NET"),
    ("Instalador",     "WiX Toolset 4 -- MSI per-user, sem necessidade de admin"),
    ("Dependencias",   "Zero Python no runtime do app -- 100% PowerShell"),
], font_size=14)
add_footer(s, 4, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Slide 5 -- Pipeline (fluxo de geracao)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_header(s, "Pipeline de geracao",
           "Da planilha ate o relatorio HTML em uma execucao")

step_w = Inches(2.3)
step_h = Inches(1.3)
gap_w = Inches(0.25)
arrow_w = Inches(0.25)
y = Inches(2.0)

# Step boxes -- 4 boxes + 3 arrows
labels = [
    ("Planilha",      "Ferias-template.xlsx\n(template oficial)"),
    ("Validacao",     "Nome + 5 abas\nobrigatorias"),
    ("Template MD",   "Substitui placeholders\n(dashboard, gantt, etc)"),
    ("Pandoc",        "MD -> HTML\nCSS + Mermaid embutidos"),
]
left = Inches(0.5)
for i, (title, body) in enumerate(labels):
    add_pipeline_step(s, left, y, step_w, step_h, title, body, color=ACCENT)
    if i < len(labels) - 1:
        add_arrow(s, left + step_w + Inches(0.02), y + Inches(0.5),
                  arrow_w, Inches(0.3))
    left = left + step_w + arrow_w + Inches(0.05)

# Output row
add_textbox(s, Inches(0.5), Inches(3.7), Inches(12), Inches(0.4),
            "Saidas geradas (sempre na mesma pasta, sobrescreve a anterior):",
            font_size=14, bold=True, color=PRIMARY)

out_w = Inches(2.85)
out_h = Inches(1.0)
out_y = Inches(4.2)
outs = [
    ("Ferias.html",  "Relatorio executivo\nestilizado"),
    ("Ferias.md",    "Versao Markdown\n(rastreabilidade)"),
    ("Ferias.xlsx",  "Copia da planilha\nque gerou o relatorio"),
    ("Ferias.pdf",   "Opcional -- compativel\ncom SharePoint"),
]
left = Inches(0.6)
for title, body in outs:
    add_pipeline_step(s, left, out_y, out_w, out_h, title, body, color=PRIMARY)
    left += out_w + Inches(0.2)

add_footer(s, 5, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Slide 6 -- GUI / experiencia do usuario
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_header(s, "Experiencia do usuario",
           "Janela unica, simples e direta")
add_bullets(s, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), [
    ("Duplo-clique no atalho 'Gerar Relatorio' -- a janela abre", True),
    "Campo Ano -- spinner com o ano atual como default.",
    "Campo Planilha -- caminho + botao 'Procurar' (filtrado para o template).",
    "Campo Pasta de saida -- caminho + 'Procurar', persistido entre execucoes (registry).",
    "Checkbox 'Abrir HTML apos gerar' -- marcado por default.",
    "Botao 'Gerar Relatorio' (atalho Enter) e 'Fechar' (atalho Esc).",
    "Barra de progresso durante a execucao -- feedback visual constante.",
    ("Mensagens de erro claras em todos os caminhos de falha (MessageBox).", True),
])
add_footer(s, 6, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Slide 7 -- Validacoes e robustez
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_header(s, "Validacoes e robustez",
           "Regras que protegem o usuario contra entradas invalidas")
add_bullets(s, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), [
    ("Instancia unica por PC", True),
    "Mutex global impede duas janelas abertas simultaneamente -- a segunda traz a existente para frente.",
    ("Lock ao template oficial", True),
    "App so aceita Ferias-template.xlsx com as 5 abas obrigatorias (Ferias, Squads, Integrantes, Status, Instrucoes).",
    "Qualquer outro arquivo e rejeitado com mensagem clara antes da leitura.",
    ("Pre-flight checks na geracao", True),
    "Verifica: planilha existe, pasta de saida existe e e gravavel, nome bate com o template.",
    ("Saida controlada", True),
    "Sempre na mesma pasta com nomes fixos -- cada execucao sobrescreve a anterior, sem lixo acumulado.",
])
add_footer(s, 7, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Slide 8 -- Distribuicao via MSI
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_header(s, "Distribuicao",
           "Instalador MSI profissional, sem necessidade de admin")
add_kv_block(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(2.4), [
    ("Instalador",  "FeriasAutomacao-1.0.0.msi  (38,5 MB)"),
    ("Tecnologia",  "WiX Toolset 4.0.6  -- per-user (LocalAppData), sem admin"),
    ("Atalhos",     "Menu Iniciar  |  Area de Trabalho  |  SendTo (clique-direito)"),
    ("Versao",      f"{VERSION}  -- mantida para esta entrega"),
], font_size=14)
add_textbox(s, Inches(0.6), Inches(4.0), Inches(12), Inches(0.4),
            "O que o usuario faz", font_size=16, bold=True, color=PRIMARY)
add_bullets(s, Inches(0.6), Inches(4.4), Inches(12), Inches(2.5), [
    "1. Recebe o .msi (e-mail, SharePoint, GitHub Release).",
    "2. Duplo-clique -- assistente Next-Next-Finish, sem prompt de admin.",
    "3. Atalho 'Gerar Relatorio' aparece no Menu Iniciar e na area de trabalho.",
    "4. Pronto. Roda igual qualquer aplicativo do Windows.",
])
add_footer(s, 8, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Slide 9 -- O que foi entregue (changelog desta entrega)
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_header(s, "O que foi feito nesta entrega",
           "Resumo dos commits mergeados na main (branch acerto-de-planilha)")
add_bullets(s, Inches(0.6), Inches(1.3), Inches(12), Inches(5.5), [
    ("Renomeacao da planilha-fonte para Ferias-template.xlsx + abas de referencia.", True),
    ("Saida controlada: pasta unica configuravel com nomes fixos, sobrescreve a anterior.", True),
    ("Single-instance lock: apenas uma janela do app por PC.", True),
    ("Validacao estrita: app so aceita a planilha-template oficial.", True),
    ("Mensagens de erro claras em todos os caminhos de falha (planilha, pasta, permissao).", True),
    ("Documentacao: troubleshooting de reparo manual do MSI no README.", True),
    ("TODO.md versionado mapeando proximas melhorias (self-healing, etc).", True),
    "Build do MSI 1.0.0 gerado e validado para distribuicao.",
])
add_footer(s, 9, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Slide 10 -- Documentacao entregue
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_header(s, "Documentacao entregue",
           "Tres niveis de documentacao para publicos diferentes")
add_kv_block(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(4.5), [
    ("README.md",
     "Visao geral tecnica + parametros CLI + troubleshooting (para devs)"),
    ("Manual do Usuario (.docx)",
     "Passo-a-passo completo: instalar -> preencher -> gerar (para usuarios finais)"),
    ("TODO.md",
     "Lista versionada de melhorias mapeadas mas nao priorizadas"),
    ("Comentarios no codigo",
     "executar.ps1 e gui.ps1 documentam regras de negocio e validacoes"),
    ("Repo GitHub",
     "github.com/ceseidl/FeriasAutomacao -- branch main atualizado"),
], font_size=14)
add_footer(s, 10, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Slide 11 -- Status atual e proximos passos
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_header(s, "Status atual e proximos passos")

# Status atual
add_textbox(s, Inches(0.6), Inches(1.2), Inches(6), Inches(0.4),
            "Status atual", font_size=18, bold=True, color=SUCCESS)
add_bullets(s, Inches(0.6), Inches(1.7), Inches(6.0), Inches(5), [
    "App em producao, versao 1.0.0",
    "Codigo na main, mergeado e empurrado",
    "MSI buildado e pronto para distribuir",
    "Documentacao completa (README + Manual + TODO)",
    "Validacoes e mensagens de erro cobrindo todos os caminhos",
], font_size=14)

# Proximos passos
add_textbox(s, Inches(7.0), Inches(1.2), Inches(6), Inches(0.4),
            "Proximos passos (mapeados)", font_size=18, bold=True, color=WARN)
add_bullets(s, Inches(7.0), Inches(1.7), Inches(6.0), Inches(5), [
    "Estudo: self-healing real do MSI (auto-reparo de arquivos deletados)",
    "Distribuicao oficial: GitHub Release / SharePoint corporativo",
    "Coletar feedback dos primeiros usuarios -- iterar conforme demanda",
    "Possivel evolucao: dashboard agregando varios meses",
], font_size=14)

add_footer(s, 11, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Slide 12 -- Numeros / metricas
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_header(s, "Numeros e metricas",
           "Dimensao do que foi entregue")

card_w = Inches(3.0)
card_h = Inches(2.0)
gap = Inches(0.2)
y1 = Inches(1.6)
y2 = Inches(3.9)

cards = [
    ("1.0.0",       "Versao oficial",          PRIMARY),
    ("38,5 MB",     "Tamanho do MSI",          PRIMARY),
    ("4 formatos",  "HTML / MD / XLSX / PDF",  PRIMARY),
    ("0",           "Comandos no CLI\npara o usuario final",      ACCENT),
    ("5 abas",      "Validadas no template\noficial obrigatorio", ACCENT),
    ("3 atalhos",   "Menu Iniciar, Desktop\ne SendTo",            ACCENT),
]
positions = [
    (Inches(0.6) + i * (card_w + gap), y1) for i in range(3)
] + [
    (Inches(0.6) + i * (card_w + gap), y2) for i in range(3)
]
for (label_top, label_bot, fill), (x, y) in zip(cards, positions):
    add_filled_rect(s, x, y, card_w, card_h, fill)
    add_textbox(s, x, y + Inches(0.25), card_w, Inches(0.9),
                label_top, font_size=42, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(s, x + Inches(0.1), y + Inches(1.25),
                card_w - Inches(0.2), Inches(0.7),
                label_bot, font_size=12, color=WHITE,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.TOP)

add_footer(s, 12, TOTAL_SLIDES)


# ===========================================================================
# "ONE MORE THING" -- estilo keynote Apple. Apos o slide de metricas a
# apresentacao parece estar terminando, mas a gente puxa uma surpresa:
# a v1.0.0.1 ja existe e tem uma feature nova (Controle de Vencimento).
# Sequencia: transicao dramatica -> reveal da feature -> comparativo visual.
# ===========================================================================

# ---------------------------------------------------------------------------
# Slide 13 -- "And one more thing..." (transicao Apple keynote)
# Fundo preto, texto branco enorme, sem header/footer. Estilo minimalista
# pra criar o efeito de surpresa antes do reveal.
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
BLACK = RGBColor(0x00, 0x00, 0x00)
add_filled_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, BLACK)
add_textbox(s, Inches(0.5), Inches(2.3), Inches(12.3), Inches(2.0),
            "And one more thing...",
            font_size=72, bold=False, color=WHITE,
            align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_textbox(s, Inches(0.5), Inches(4.3), Inches(12.3), Inches(0.6),
            "(ainda nao acabou)",
            font_size=18, color=RGBColor(0x88, 0x88, 0x88),
            align=PP_ALIGN.CENTER)


# ---------------------------------------------------------------------------
# Slide 14 -- Reveal da v1.0.0.1: Controle de Vencimento de Ferias
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_header(s, "v1.0.0.1  --  Controle de Vencimento de Ferias",
           "Em desenvolvimento -- ja temos um build funcionando")

# Cartao de destaque com a feature nova
add_filled_rect(s, Inches(0.6), Inches(1.4), Inches(12.1), Inches(1.3), ACCENT)
add_textbox(s, Inches(0.8), Inches(1.55), Inches(11.7), Inches(0.45),
            "Nova secao no relatorio: alerta automatico de quem precisa tirar ferias",
            font_size=18, bold=True, color=WHITE)
add_textbox(s, Inches(0.8), Inches(2.05), Inches(11.7), Inches(0.65),
            "Calcula 1o e 2o vencimento por colaborador (CLT 12/24 meses) a partir da data de admissao "
            "e da data fim das ultimas ferias na aba Integrantes.",
            font_size=13, color=WHITE)

# Bloco com a regra de negocio
add_textbox(s, Inches(0.6), Inches(2.95), Inches(12), Inches(0.4),
            "3 sub-blocos automaticos no relatorio:",
            font_size=15, bold=True, color=PRIMARY)

block_w = Inches(4.0)
block_h = Inches(2.4)
block_y = Inches(3.45)
block_gap = Inches(0.07)

# CRITICO
x = Inches(0.6)
add_filled_rect(s, x, block_y, block_w, block_h, RGBColor(0xE5, 0x3E, 0x3E))
add_textbox(s, x + Inches(0.2), block_y + Inches(0.15), block_w - Inches(0.4), Inches(0.45),
            "CRITICO -- 2o vencimento", font_size=15, bold=True, color=WHITE)
add_textbox(s, x + Inches(0.2), block_y + Inches(0.65), block_w - Inches(0.4), block_h - Inches(0.7),
            "Pessoas com ferias proximas de vencer EM DOBRO ou ja vencidas.\n\n"
            "Janela de alerta: 6 meses antes do 2o vencimento.\n\n"
            "Acao: tirar ferias urgente.",
            font_size=11, color=WHITE)

# ATENCAO
x = Inches(0.6) + block_w + block_gap
add_filled_rect(s, x, block_y, block_w, block_h, RGBColor(0xED, 0x89, 0x36))
add_textbox(s, x + Inches(0.2), block_y + Inches(0.15), block_w - Inches(0.4), Inches(0.45),
            "ATENCAO -- 1o vencimento", font_size=15, bold=True, color=WHITE)
add_textbox(s, x + Inches(0.2), block_y + Inches(0.65), block_w - Inches(0.4), block_h - Inches(0.7),
            "Pessoas que ja adquiriram o direito ou estao perto.\n\n"
            "Janela de alerta: 6 meses antes do 1o vencimento.\n\n"
            "Acao: comecar a planejar ferias.",
            font_size=11, color=WHITE)

# DADOS INCOMPLETOS
x = Inches(0.6) + 2 * (block_w + block_gap)
add_filled_rect(s, x, block_y, block_w, block_h, RGBColor(0xA0, 0xAE, 0xC0))
add_textbox(s, x + Inches(0.2), block_y + Inches(0.15), block_w - Inches(0.4), Inches(0.45),
            "Dados incompletos", font_size=15, bold=True, color=WHITE)
add_textbox(s, x + Inches(0.2), block_y + Inches(0.65), block_w - Inches(0.4), block_h - Inches(0.7),
            "Pessoas sem cadastro completo na aba Integrantes.\n\n"
            "Lista as colunas que faltam (data de admissao).\n\n"
            "Acao: atualizar a planilha.",
            font_size=11, color=WHITE)

# Linha de fechamento
add_textbox(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.5),
            "Quem esta em dia (>6 meses pra qualquer vencimento) nao aparece -- mantem a secao limpa.",
            font_size=12, color=MUTED, align=PP_ALIGN.CENTER)

add_footer(s, 14, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Slide 15 -- Comparativo visual antes / depois
# Insere a imagem comparativo-1.0.0-vs-1.0.0.1.png centralizada.
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_header(s, "Antes  vs  Depois",
           "v1.0.0 (em producao)  ->  v1.0.0.1 (em desenvolvimento)")

comparativo_path = Path(__file__).parent.parent / 'results' / 'comparativo-1.0.0-vs-1.0.0.1.png'
if comparativo_path.exists():
    # A imagem tem proporcao alta (mais alta que larga em side-by-side com
    # paginas de relatorio). Encaixa centralizada respeitando area util.
    from PIL import Image as _PILImage
    with _PILImage.open(comparativo_path) as _im:
        img_w, img_h = _im.size
    aspect = img_w / img_h

    avail_w = SLIDE_W - Inches(1.0)        # margem 0.5in cada lado
    avail_h = SLIDE_H - Inches(1.5)        # respeita header (0.95) + folga
    avail_aspect = avail_w / avail_h

    if aspect > avail_aspect:
        # imagem mais larga proporcionalmente -> limita pela largura
        place_w = avail_w
        place_h = int(avail_w / aspect)
    else:
        # imagem mais alta -> limita pela altura
        place_h = avail_h
        place_w = int(avail_h * aspect)

    place_x = Inches(0.5) + (avail_w - place_w) / 2
    place_y = Inches(1.05) + (avail_h - place_h) / 2
    s.shapes.add_picture(str(comparativo_path), place_x, place_y, width=place_w, height=place_h)
else:
    add_textbox(s, Inches(0.6), Inches(2.0), Inches(12), Inches(1.0),
                f"AVISO: imagem nao encontrada em {comparativo_path}",
                font_size=18, color=WARN, align=PP_ALIGN.CENTER)
    add_textbox(s, Inches(0.6), Inches(3.0), Inches(12), Inches(0.6),
                "Rode docs/compose-comparativo.py antes de gerar a apresentacao.",
                font_size=14, color=MUTED, align=PP_ALIGN.CENTER)

add_footer(s, 15, TOTAL_SLIDES)


# ---------------------------------------------------------------------------
# Slide 16 -- Encerramento
# ---------------------------------------------------------------------------
s = prs.slides.add_slide(blank)
add_filled_rect(s, Inches(0), Inches(0), SLIDE_W, SLIDE_H, PRIMARY)
add_filled_rect(s, Inches(0), Inches(5.4), SLIDE_W, Inches(0.06), ACCENT)
add_textbox(s, Inches(0.8), Inches(2.5), Inches(11.5), Inches(1.2),
            "Obrigado", font_size=72, bold=True, color=WHITE)
add_textbox(s, Inches(0.8), Inches(3.7), Inches(11.5), Inches(0.6),
            "Duvidas, demonstracao ao vivo ou proximos passos",
            font_size=22, color=WHITE)
add_textbox(s, Inches(0.8), Inches(5.8), Inches(11.5), Inches(0.5),
            f"{AUTHOR}", font_size=16, color=WHITE)
add_textbox(s, Inches(0.8), Inches(6.2), Inches(11.5), Inches(0.5),
            "github.com/ceseidl/FeriasAutomacao",
            font_size=14, color=WHITE)


# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
out_path = Path(__file__).parent / "Apresentacao-FeriasAutomacao.pptx"
prs.save(out_path)
print(f"OK: {out_path}  ({out_path.stat().st_size // 1024} KB)")
